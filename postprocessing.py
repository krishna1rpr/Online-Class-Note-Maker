import os
import cv2
import openai
import whisper
import subprocess
from math import log
from pptx.util import Inches
from pptx import Presentation


# utility functions

def convert_to_grayscale(image):
    return cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)


def calc_diff(i1, i2):
    i1 = convert_to_grayscale(i1)
    i2 = convert_to_grayscale(i2)

    diff = cv2.absdiff(i1, i2)

    mse = diff.mean()**0.5

    return mse


# main class

class Proj:
    def __init__(self, src, frames, diff, audio, video, pres):
        self.src = os.path.join(os.getcwd(), src)
        self.frames = os.path.join(os.getcwd(), frames)
        self.diff = os.path.join(os.getcwd(), diff)
        self.audio = os.path.join(os.getcwd(), audio)
        self.video = os.path.join(os.getcwd(), video)
        self.pres = os.path.join(os.getcwd(), pres)
        self.text = ""

        if not os.path.exists(self.frames):
            os.mkdir(self.frames)
        if not os.path.exists(self.diff):
            os.mkdir(self.diff)

    def generate_frames(self):
        if not os.path.exists(self.frames):
            os.mkdir(self.frames)
        subprocess.call(["ffmpeg", "-i", self.src, "-vf",
                        "fps=1/3.5", self.frames + "/o%03d.jpg"])

    def generate_slides(self):
        prev = None

        for file in os.listdir(self.frames):
            if file == "o001.jpg":
                prev = cv2.imread(self.frames + file)
                pass
            else:
                curr = cv2.imread(self.frames + file)
                diff = calc_diff(prev, curr)
                prev = curr
                if diff > 2.72:
                    cv2.imwrite(self.diff + file, curr)

    def segregate_audio_and_video(self):
        subprocess.call(["ffmpeg", "-i", self.src, "-vn", "-ac", "1", "-ar", "16000", "-acodec", "pcm_s16le",
                        self.audio])
        # subprocess.call(["ffmpeg", "-i", self.src, "-vcodec", "copy", "-an",
        #                 self.video])

    def generate_ppt(self):
        prs = Presentation()
        frm, too = 0, 0

        for image_file in os.listdir(self.diff):
            too = int(image_file.split(".")[0][1:])
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            left = Inches(0.5)
            top = Inches(0.5)
            height = Inches(5)
            slide.shapes.add_picture(os.path.join(
                self.diff, image_file), left, top, height=height)

            fname = f"output{too}.wav"

            subprocess.call(["ffmpeg", "-y", "-i", self.audio, "-ss",
                            str(frm*3.5), "-to", str(too*3.5), "-c", "copy", fname])

            textbox = slide.shapes.add_textbox(left=Inches(
                0.4), top=Inches(5.75), width=Inches(9), height=Inches(2))
            textbox.text = self.generate_text(fname)
            textbox.text_frame.word_wrap = True

            frm, too = too, 0

        prs.save(self.pres)

    def generate_text(self, src=None):
        model = whisper.load_model("tiny")
        result = model.transcribe(src or self.audio)

        with open("transcript.txt", "a") as f:
            f.write(result['text']+"\n")
            self.text += result['text']
            return result['text']

    def generate_summary(self):
        openai.api_key = "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"

        response = openai.Completion.create(
            model="text-davinci-003",
            prompt=self.text,
            temperature=0.7,
            max_tokens=500,
            top_p=1.0,
            frequency_penalty=0.0,
            presence_penalty=1
        )

        with open("summary.txt", "w") as f:
            f.write(response['choices'][0]['text'])


if __name__ == "__main__":
    proj = Proj("./src2.mp4", "./frames/", "./diff/",
                "./audio.wav", "./video.mp4", "./presentation.pptx")
    proj.generate_frames()
    proj.generate_slides()
    # proj.segregate_audio_and_video()
    proj.generate_ppt()
    # proj.generate_text()
    # proj.generate_summary()
